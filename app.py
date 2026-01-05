"""
Best Pdf Converter - A SaaS PDF manipulation tool
Version 2.0.0 with full authentication and database support
"""
import os
import zipfile
import secrets
import logging
from logging.handlers import RotatingFileHandler
from datetime import datetime
from flask import Flask, render_template, request, redirect, url_for, send_file, flash, jsonify, session
from flask_login import LoginManager, login_required, current_user
from flask_limiter import Limiter
from flask_limiter.util import get_remote_address
from flask_talisman import Talisman
from PyPDF2 import PdfReader, PdfWriter
from werkzeug.utils import secure_filename
from werkzeug.exceptions import RequestEntityTooLarge
from dotenv import load_dotenv

# Optional PDF2Word conversion support using PyMuPDF and python-docx
try:
    import fitz  # PyMuPDF
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    PDF2WORD_AVAILABLE = True
except ImportError:
    PDF2WORD_AVAILABLE = False

# Optional PDF page preview support
try:
    from pdf2image import convert_from_path
    from io import BytesIO
    PDF_PREVIEW_AVAILABLE = True
except ImportError:
    PDF_PREVIEW_AVAILABLE = False

# Load environment variables
load_dotenv()

# Import local modules
from config import config
from models import db, bcrypt, User, UsageRecord
from auth import auth_bp
from payment import payment_bp
from analytics import init_analytics
from analytics_middleware import setup_analytics_middleware

__version__ = "2.0.0"

def create_app(config_name=None):
    """Application factory"""
    if config_name is None:
        config_name = os.environ.get('FLASK_ENV', 'development')

    app = Flask(__name__)
    app.config.from_object(config[config_name])

    # Configure logging
    if not app.debug:
        if not os.path.exists('logs'):
            os.mkdir('logs')
        file_handler = RotatingFileHandler('logs/zenpdf.log', maxBytes=10240000, backupCount=10)
        file_handler.setFormatter(logging.Formatter(
            '%(asctime)s %(levelname)s: %(message)s [in %(pathname)s:%(lineno)d]'
        ))
        file_handler.setLevel(logging.INFO)
        app.logger.addHandler(file_handler)
        app.logger.setLevel(logging.INFO)
        app.logger.info('ZenPDF startup')

    # Log configuration warnings (avoid reentrant logging issues)
    if config_name == 'production':
        if not os.environ.get('DATABASE_URL'):
            app.logger.warning(
                "DATABASE_URL not set - using SQLite. "
                "For production, add PostgreSQL database in Railway dashboard."
            )
        if not os.environ.get('SECRET_KEY'):
            app.logger.warning(
                "SECRET_KEY not set - using default (INSECURE!). "
                "Set SECRET_KEY environment variable in Railway."
            )

    # Initialize extensions
    db.init_app(app)
    bcrypt.init_app(app)

    # Flask-Login setup
    login_manager = LoginManager()
    login_manager.init_app(app)
    login_manager.login_view = 'auth.login'
    login_manager.login_message = 'Please log in to access this page.'
    login_manager.login_message_category = 'info'

    @login_manager.user_loader
    def load_user(user_id):
        return User.query.get(int(user_id))

    # Rate limiting
    limiter = Limiter(
        app=app,
        key_func=get_remote_address,
        default_limits=["200 per day", "50 per hour"],
        storage_uri=app.config['RATELIMIT_STORAGE_URL']
    )

    # Security headers (disable in development)
    if not app.debug:
        Talisman(
            app,
            force_https=app.config['TALISMAN_FORCE_HTTPS'],
            strict_transport_security=app.config['TALISMAN_STRICT_TRANSPORT_SECURITY'],
            content_security_policy=app.config['TALISMAN_CONTENT_SECURITY_POLICY']
        )

    # Register blueprints
    app.register_blueprint(auth_bp, url_prefix='/auth')
    app.register_blueprint(payment_bp, url_prefix='/payment')

    # Initialize server-side analytics
    from analytics import analytics
    init_analytics(app)
    if analytics:
        setup_analytics_middleware(app, analytics)

    # Context processor to inject Google Drive config into templates
    @app.context_processor
    def inject_google_config():
        return {
            'google_config': {
                'api_key': os.environ.get('GOOGLE_PICKER_API_KEY', ''),
                'client_id': os.environ.get('GOOGLE_DRIVE_CLIENT_ID', ''),
                'app_id': os.environ.get('GOOGLE_APP_ID', '')
            }
        }

    # ========== HELPER FUNCTIONS ==========

    def cleanup_old_files(folder, max_age_hours=1):
        """Remove files older than max_age_hours"""
        try:
            if not os.path.exists(folder):
                return
            current_time = datetime.now().timestamp()
            for filename in os.listdir(folder):
                file_path = os.path.join(folder, filename)
                if os.path.isfile(file_path):
                    file_age = current_time - os.path.getmtime(file_path)
                    if file_age > (max_age_hours * 3600):
                        try:
                            os.remove(file_path)
                            app.logger.info(f'Cleaned up old file: {filename}')
                        except OSError as e:
                            app.logger.error(f'Failed to remove old file {filename}: {e}')
        except Exception as e:
            app.logger.error(f'Error during file cleanup in {folder}: {e}')

    def allowed_file(filename):
        """Check if file extension is allowed"""
        return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

    def sanitize_filename(filename):
        """Sanitize filename to prevent path traversal"""
        filename = os.path.basename(filename)
        filename = secure_filename(filename)
        if not filename.lower().endswith('.pdf'):
            filename += '.pdf'
        return filename

    def validate_file_size(file):
        """Validate file size based on user tier"""
        file.seek(0, os.SEEK_END)
        size = file.tell()
        file.seek(0)

        # Premium users get higher limits, free users and anonymous get standard limit
        is_premium = current_user.is_authenticated and current_user.is_premium
        max_size = app.config['MAX_FILE_SIZE_PREMIUM'] if is_premium else app.config['MAX_FILE_SIZE']

        if size > max_size:
            max_mb = max_size // (1024 * 1024)
            raise ValueError(f"File size exceeds {max_mb}MB limit")
        if size == 0:
            raise ValueError("File is empty")
        return size

    def record_usage(operation_type, file_size=None, pages_processed=None, success=True, error_message=None):
        """Record usage for analytics and rate limiting"""
        if current_user.is_authenticated:
            usage = UsageRecord(
                user_id=current_user.id,
                operation_type=operation_type,
                file_size=file_size,
                pages_processed=pages_processed,
                ip_address=request.remote_addr,
                user_agent=request.headers.get('User-Agent', '')[:255],
                success=success,
                error_message=error_message
            )
            db.session.add(usage)
            db.session.commit()

    def check_usage_limit():
        """Check if user has exceeded their usage limit - BASIC TOOLS ARE UNLIMITED"""
        # All basic PDF operations are FREE and UNLIMITED for everyone
        return True, None
    
    def check_ai_usage_limit():
        """Check if user can use AI features based on their credits"""
        if not current_user.is_authenticated:
            # Anonymous users cannot use AI features - require login
            return False, "Please sign up or log in to use AI features. Get 3 free AI credits every month!"
        
        if current_user.can_use_ai_feature():
            remaining = current_user.get_ai_credits_remaining()
            limit = current_user.get_ai_credits_limit()
            return True, f"AI credits: {remaining}/{limit} remaining this month"
        else:
            return False, f"You've used all your AI credits this month. Upgrade to Pro for 20 credits/month or Business for unlimited!"

    # Ensure upload folders exist
    for folder in [app.config['UPLOAD_FOLDER'], app.config['SPLIT_FOLDER'], app.config['MERGED_FOLDER']]:
        os.makedirs(folder, exist_ok=True)

    # Clean up old files on startup
    cleanup_old_files(app.config['SPLIT_FOLDER'], app.config['FILE_CLEANUP_HOURS'])
    cleanup_old_files(app.config['MERGED_FOLDER'], app.config['FILE_CLEANUP_HOURS'])
    cleanup_old_files(app.config['UPLOAD_FOLDER'], app.config['FILE_CLEANUP_HOURS'])

    # ========== ERROR HANDLERS ==========

    @app.errorhandler(413)
    @app.errorhandler(RequestEntityTooLarge)
    def too_large(e):
        app.logger.warning(f'File too large from {request.remote_addr}')
        flash('File too large. Maximum size is 50MB for free users, 100MB for premium.', 'error')
        return redirect(request.url), 413

    @app.errorhandler(500)
    def internal_error(e):
        app.logger.error(f'Internal error: {e}')
        db.session.rollback()
        flash('An error occurred. Please try again.', 'error')
        return redirect(url_for('home')), 500

    @app.errorhandler(404)
    def not_found(e):
        return render_template('index.html'), 404

    # ========== PUBLIC ROUTES ==========

    @app.route('/')
    def home():
        return render_template('index.html', version=__version__)

    @app.route('/sitemap.xml')
    def sitemap():
        return send_file('static/sitemap.xml', mimetype='application/xml')

    @app.route('/robots.txt')
    def robots():
        return send_file('static/robots.txt', mimetype='text/plain')

    @app.route('/version')
    def version():
        """API endpoint to check version"""
        return jsonify({
            'version': __version__,
            'name': 'ZenPDF',
            'status': 'production'
        })

    @app.route('/health')
    def health():
        """Health check endpoint for monitoring"""
        try:
            # Check database connection
            db.session.execute(db.text('SELECT 1'))
            db_status = 'healthy'
        except Exception as e:
            app.logger.error(f'Database health check failed: {e}')
            db_status = 'unhealthy'

        return jsonify({
            'status': 'healthy' if db_status == 'healthy' else 'degraded',
            'version': __version__,
            'database': db_status,
            'timestamp': datetime.utcnow().isoformat()
        }), 200 if db_status == 'healthy' else 503

    @app.route('/api/live-stats')
    def live_stats():
        """API endpoint for live activity stats on homepage"""
        try:
            today_start = datetime.utcnow().replace(hour=0, minute=0, second=0, microsecond=0)
            today_operations = UsageRecord.query.filter(UsageRecord.created_at >= today_start).count()
            total_operations = UsageRecord.query.count()
            total_users = User.query.count()
            
            # Add some padding to make numbers look more impressive initially
            # These will naturally grow as real users come in
            display_today = today_operations + 47  # baseline activity
            display_total = total_operations + 5000  # baseline total
            display_users = total_users + 150  # baseline users
            
            return jsonify({
                'today_operations': display_today,
                'total_operations': display_total,
                'total_users': display_users,
                'timestamp': datetime.utcnow().isoformat()
            })
        except Exception as e:
            app.logger.error(f'Live stats error: {e}')
            return jsonify({
                'today_operations': 47,
                'total_operations': 5000,
                'total_users': 150,
                'timestamp': datetime.utcnow().isoformat()
            })

    @app.route('/pricing')
    def pricing():
        """Pricing page"""
        return render_template('pricing.html')

    @app.route('/privacy')
    def privacy_policy():
        """Privacy Policy page"""
        return render_template('privacy.html')

    @app.route('/terms')
    def terms_of_service():
        """Terms of Service page"""
        return render_template('terms.html')

    # ========== BLOG ==========

    @app.route('/blog')
    def blog():
        """Blog listing page"""
        return render_template('blog.html')

    @app.route('/blog/how-to-compress-pdf-without-losing-quality')
    def blog_compress():
        """Blog post: How to compress PDF"""
        return render_template('blog/how-to-compress-pdf-without-losing-quality.html')

    @app.route('/blog/merge-pdf-files-free-online')
    def blog_merge():
        """Blog post: How to merge PDF files"""
        return render_template('blog/merge-pdf-files-free-online.html')

    @app.route('/blog/split-pdf-extract-pages')
    def blog_split():
        """Blog post: How to split PDF"""
        return render_template('blog/split-pdf-extract-pages.html')

    @app.route('/blog/convert-pdf-to-word-editable')
    def blog_pdf2word():
        """Blog post: Convert PDF to Word"""
        return render_template('blog/convert-pdf-to-word-editable.html')

    @app.route('/blog/ocr-pdf-extract-text-scanned')
    def blog_ocr():
        """Blog post: OCR PDF extract text"""
        return render_template('blog/ocr-pdf-extract-text-scanned.html')

    @app.route('/blog/rotate-pdf-fix-orientation')
    def blog_rotate():
        """Blog post: Rotate PDF fix orientation"""
        return render_template('blog/rotate-pdf-fix-orientation.html')

    # ========== DASHBOARD ==========

    @app.route('/dashboard')
    @login_required
    def dashboard():
        """User dashboard"""
        # Get usage stats
        daily_usage = current_user.get_daily_usage_count()
        monthly_usage = current_user.get_monthly_usage_count()

        # Get recent operations
        recent_operations = current_user.usage_records.order_by(
            UsageRecord.created_at.desc()
        ).limit(10).all()

        return render_template('dashboard.html',
                             daily_usage=daily_usage,
                             monthly_usage=monthly_usage,
                             recent_operations=recent_operations)

    # ========== ADMIN PANEL ==========

    @app.route('/admin')
    @login_required
    def admin_panel():
        """Admin panel - only accessible to admin users"""
        # Check if user is admin
        if not current_user.is_admin:
            flash('Access denied. Admin privileges required.', 'error')
            return redirect(url_for('dashboard'))

        # Get all users
        users = User.query.order_by(User.created_at.desc()).all()

        # Get system statistics
        total_users = User.query.count()
        active_users = User.query.filter_by(is_active=True).count()
        premium_users = User.query.filter(User.subscription_tier != 'free').count()

        # Get all operations
        total_operations = UsageRecord.query.count()
        today_start = datetime.utcnow().replace(hour=0, minute=0, second=0, microsecond=0)
        today_operations = UsageRecord.query.filter(UsageRecord.created_at >= today_start).count()

        # Get recent operations across all users
        recent_operations = UsageRecord.query.order_by(
            UsageRecord.created_at.desc()
        ).limit(20).all()

        return render_template('admin_panel.html',
                             users=users,
                             total_users=total_users,
                             active_users=active_users,
                             premium_users=premium_users,
                             total_operations=total_operations,
                             today_operations=today_operations,
                             recent_operations=recent_operations)

    @app.route('/admin/user/<int:user_id>/toggle-active', methods=['POST'])
    @login_required
    def admin_toggle_user_active(user_id):
        """Toggle user active status"""
        if not current_user.is_admin:
            return jsonify({'error': 'Access denied'}), 403

        user = User.query.get_or_404(user_id)
        user.is_active = not user.is_active
        db.session.commit()

        status = 'activated' if user.is_active else 'deactivated'
        flash(f'User {user.email} has been {status}.', 'success')
        return redirect(url_for('admin_panel'))

    @app.route('/admin/user/<int:user_id>/update-tier', methods=['POST'])
    @login_required
    def admin_update_user_tier(user_id):
        """Update user subscription tier"""
        if not current_user.is_admin:
            return jsonify({'error': 'Access denied'}), 403

        user = User.query.get_or_404(user_id)
        new_tier = request.form.get('tier')

        if new_tier not in ['free', 'premium', 'enterprise']:
            flash('Invalid subscription tier.', 'error')
            return redirect(url_for('admin_panel'))

        user.subscription_tier = new_tier

        # If upgrading to premium/enterprise, set subscription dates
        if new_tier != 'free':
            user.subscription_start = datetime.utcnow()
            # Set to 1 year from now
            from datetime import timedelta
            user.subscription_end = datetime.utcnow() + timedelta(days=365)
        else:
            user.subscription_start = None
            user.subscription_end = None

        db.session.commit()
        flash(f'User {user.email} subscription updated to {new_tier}.', 'success')
        return redirect(url_for('admin_panel'))

    @app.route('/admin/user/<int:user_id>/toggle-admin', methods=['POST'])
    @login_required
    def admin_toggle_user_admin(user_id):
        """Toggle user admin status"""
        if not current_user.is_admin:
            return jsonify({'error': 'Access denied'}), 403

        user = User.query.get_or_404(user_id)

        # Prevent removing admin from yourself
        if user.id == current_user.id:
            flash('You cannot remove admin privileges from yourself.', 'error')
            return redirect(url_for('admin_panel'))

        user.is_admin = not user.is_admin
        db.session.commit()

        status = 'granted' if user.is_admin else 'revoked'
        flash(f'Admin privileges {status} for {user.email}.', 'success')
        return redirect(url_for('admin_panel'))

    # ========== DOWNLOAD ENDPOINT ==========

    @app.route('/download/<filename>')
    def download_file(filename):
        """Download processed PDF file"""
        # Check in all possible folders
        for folder in [app.config['SPLIT_FOLDER'], app.config['MERGED_FOLDER'], app.config['UPLOAD_FOLDER']]:
            file_path = os.path.join(folder, filename)
            if os.path.exists(file_path):
                return send_file(file_path, as_attachment=True, download_name=filename, mimetype='application/pdf' if filename.endswith('.pdf') else 'application/zip')

        flash('File not found or has expired.', 'error')
        return redirect(url_for('home'))

    # ========== PDF PAGE PREVIEW ==========

    @app.route('/api/pdf-preview/<filename>/<int:page_num>')
    def pdf_page_preview(filename, page_num):
        """Generate a thumbnail preview of a specific PDF page"""
        if not PDF_PREVIEW_AVAILABLE:
            return jsonify({'error': 'Preview not available'}), 503
        
        # Find the file
        file_path = None
        for folder in [app.config['UPLOAD_FOLDER'], app.config['MERGED_FOLDER']]:
            potential_path = os.path.join(folder, filename)
            if os.path.exists(potential_path):
                file_path = potential_path
                break
        
        if not file_path:
            return jsonify({'error': 'File not found'}), 404
        
        try:
            # Convert specific page to image (page_num is 1-indexed)
            images = convert_from_path(
                file_path,
                first_page=page_num,
                last_page=page_num,
                size=(150, None),  # 150px width, maintain aspect ratio
                fmt='jpeg'
            )
            
            if not images:
                return jsonify({'error': 'Could not generate preview'}), 500
            
            # Convert to bytes
            img_byte_arr = BytesIO()
            images[0].save(img_byte_arr, format='JPEG', quality=70)
            img_byte_arr.seek(0)
            
            return send_file(
                img_byte_arr,
                mimetype='image/jpeg',
                as_attachment=False
            )
        except Exception as e:
            app.logger.error(f'Preview generation error: {e}')
            return jsonify({'error': str(e)}), 500

    # ========== PDF OPERATIONS ==========

    @app.route('/split', methods=['GET', 'POST'])
    @limiter.limit("30 per hour")
    def split_pdf():
        if request.method == 'POST':
            if 'file' in request.files:
                # Check usage limit
                can_proceed, error_msg = check_usage_limit()
                if not can_proceed:
                    flash(error_msg, 'error')
                    return render_template('split.html', file_uploaded=False)

                file = request.files['file']
                if file and file.filename and allowed_file(file.filename):
                    try:
                        file_size = validate_file_size(file)

                        safe_filename = sanitize_filename(file.filename)
                        unique_filename = f"{secrets.token_hex(8)}_{safe_filename}"
                        file_path = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)
                        file.save(file_path)

                        try:
                            reader = PdfReader(file_path)
                            page_count = len(reader.pages)
                            if page_count == 0:
                                os.remove(file_path)
                                flash('Invalid PDF: File has no pages.', 'error')
                                return render_template('split.html', file_uploaded=False)
                        except Exception as e:
                            os.remove(file_path)
                            flash(f'Invalid PDF file: {str(e)}', 'error')
                            return render_template('split.html', file_uploaded=False)

                        session['split_file'] = unique_filename
                        flash('PDF uploaded successfully!', 'success')
                        return render_template('split.html', file_uploaded=True, page_count=page_count, file_name=unique_filename)

                    except ValueError as e:
                        flash(str(e), 'error')
                    except Exception as e:
                        flash(f'Error uploading file: {str(e)}', 'error')
                else:
                    flash('Please select a valid PDF file.', 'error')

            elif 'split_type' in request.form:
                try:
                    file_name = session.get('split_file') or request.form.get('file_name')
                    if not file_name:
                        flash('File not found. Please upload again.', 'error')
                        return render_template('split.html', file_uploaded=False)

                    file_path = os.path.join(app.config['UPLOAD_FOLDER'], file_name)

                    if not os.path.exists(file_path):
                        flash('File not found. Please upload again.', 'error')
                        session.pop('split_file', None)
                        return render_template('split.html', file_uploaded=False)

                    split_type = request.form['split_type']
                    pages_to_split = []

                    reader = PdfReader(file_path)
                    total_pages = len(reader.pages)

                    if split_type == 'range':
                        try:
                            page_ranges = request.form.get('page_ranges', '').strip()
                            if not page_ranges:
                                flash('Please enter page ranges.', 'error')
                                return render_template('split.html', file_uploaded=True,
                                                     page_count=total_pages, file_name=file_name)
                            
                            # Parse ranges - keep them as separate ranges, don't flatten
                            ranges_list = []  # List of (start, end) tuples (0-indexed)
                            for part in page_ranges.split(','):
                                part = part.strip()
                                if '-' in part:
                                    start, end = map(int, part.split('-'))
                                    if start > end:
                                        flash(f'Invalid range: {start}-{end}. Start must be less than or equal to end.', 'error')
                                        return render_template('split.html', file_uploaded=True,
                                                             page_count=total_pages, file_name=file_name)
                                    if start < 1 or end > total_pages:
                                        flash(f'Invalid range: {start}-{end}. Pages must be between 1 and {total_pages}.', 'error')
                                        return render_template('split.html', file_uploaded=True,
                                                             page_count=total_pages, file_name=file_name)
                                    ranges_list.append((start - 1, end - 1))  # Convert to 0-indexed
                                else:
                                    page = int(part)
                                    if page < 1 or page > total_pages:
                                        flash(f'Invalid page: {page}. Pages must be between 1 and {total_pages}.', 'error')
                                        return render_template('split.html', file_uploaded=True,
                                                             page_count=total_pages, file_name=file_name)
                                    ranges_list.append((page - 1, page - 1))  # Single page as range
                            
                            if not ranges_list:
                                flash('No valid ranges specified.', 'error')
                                return render_template('split.html', file_uploaded=True,
                                                     page_count=total_pages, file_name=file_name)
                            
                            # Use the new range-based split function
                            zip_file_path = split_pdf_by_ranges(file_path, ranges_list, app.config['SPLIT_FOLDER'])
                            zip_filename = os.path.basename(zip_file_path)
                            
                            # Count total pages for stats
                            total_pages_extracted = sum(end - start + 1 for start, end in ranges_list)
                            
                            # Record usage
                            file_size = os.path.getsize(file_path)
                            record_usage('split', file_size=file_size, pages_processed=total_pages_extracted)

                            # Clean up uploaded file
                            try:
                                os.remove(file_path)
                                session.pop('split_file', None)
                            except:
                                pass

                            # Return success page with download button
                            return render_template('success.html',
                                                 operation='Split',
                                                 filename=zip_filename,
                                                 pages_count=total_pages_extracted)
                        except ValueError:
                            flash('Invalid page range format. Use ranges like 1-5, 10-15.', 'error')
                            return render_template('split.html', file_uploaded=True,
                                                 page_count=total_pages, file_name=file_name)

                    elif split_type == 'specific':
                        try:
                            specific_pages = request.form.get('specific_pages', '').strip()
                            page_set = set()
                            for part in specific_pages.split(','):
                                part = part.strip()
                                if '-' in part:
                                    start, end = map(int, part.split('-'))
                                    page_set.update(range(start, end + 1))
                                else:
                                    page_set.add(int(part))

                            pages_to_split = sorted([p - 1 for p in page_set if 1 <= p <= total_pages])

                            if not pages_to_split:
                                flash('No valid pages specified.', 'error')
                                return render_template('split.html', file_uploaded=True,
                                                     page_count=total_pages, file_name=file_name)
                        except ValueError:
                            flash('Invalid page format. Use comma-separated numbers or ranges (e.g., 1,3,5-7).', 'error')
                            return render_template('split.html', file_uploaded=True,
                                                 page_count=total_pages, file_name=file_name)

                    elif split_type == 'even':
                        pages_to_split = [i for i in range(total_pages) if (i + 1) % 2 == 0]
                    elif split_type == 'odd':
                        pages_to_split = [i for i in range(total_pages) if (i + 1) % 2 != 0]
                    else:
                        flash('Invalid split type.', 'error')
                        return render_template('split.html', file_uploaded=True,
                                             page_count=total_pages, file_name=file_name)

                    if not pages_to_split:
                        flash('No pages selected to split.', 'error')
                        return render_template('split.html', file_uploaded=True,
                                             page_count=total_pages, file_name=file_name)

                    # Split PDF and create ZIP
                    zip_file_path = split_pdf_pages(file_path, pages_to_split, app.config['SPLIT_FOLDER'])
                    zip_filename = os.path.basename(zip_file_path)

                    # Record usage
                    file_size = os.path.getsize(file_path)
                    record_usage('split', file_size=file_size, pages_processed=len(pages_to_split))

                    # Clean up uploaded file
                    try:
                        os.remove(file_path)
                        session.pop('split_file', None)
                    except:
                        pass

                    # Return success page with download button instead of auto-download
                    return render_template('success.html',
                                         operation='Split',
                                         filename=zip_filename,
                                         pages_count=len(pages_to_split))

                except Exception as e:
                    record_usage('split', success=False, error_message=str(e))
                    flash(f'Error processing PDF: {str(e)}', 'error')
                    session.pop('split_file', None)
                    return render_template('split.html', file_uploaded=False)

        return render_template('split.html', file_uploaded=False)

    # Similar refactoring needed for compress, rotate, and merge routes
    # (keeping original logic but adding authentication and usage tracking)

    @app.route('/compress', methods=['GET', 'POST'])
    @limiter.limit("30 per hour")
    def compress_pdf():
        if request.method == 'POST':
            if 'file' not in request.files:
                flash('No file selected.', 'error')
                return render_template('compress.html', file_uploaded=False)

            # Check usage limit
            can_proceed, error_msg = check_usage_limit()
            if not can_proceed:
                flash(error_msg, 'error')
                return render_template('compress.html', file_uploaded=False)

            file = request.files['file']
            if file and file.filename and allowed_file(file.filename):
                try:
                    file_size = validate_file_size(file)

                    safe_filename = sanitize_filename(file.filename)
                    unique_filename = f"{secrets.token_hex(8)}_{safe_filename}"
                    file_path = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)
                    file.save(file_path)

                    try:
                        reader = PdfReader(file_path)
                        page_count = len(reader.pages)
                        if page_count == 0:
                            os.remove(file_path)
                            flash('Invalid PDF: File has no pages.', 'error')
                            return render_template('compress.html', file_uploaded=False)
                    except Exception as e:
                        if os.path.exists(file_path):
                            os.remove(file_path)
                        app.logger.error(f'PDF validation error: {e}')
                        flash(f'Invalid PDF file: {str(e)}', 'error')
                        return render_template('compress.html', file_uploaded=False)

                    # Compress the PDF
                    compressed_path = compress_pdf_file(file_path)
                    compressed_filename = os.path.basename(compressed_path)

                    # Record usage
                    record_usage('compress', file_size=file_size, pages_processed=page_count)

                    # Clean up original file
                    try:
                        os.remove(file_path)
                    except OSError as e:
                        app.logger.error(f'Failed to remove temp file: {e}')

                    # Return success page with download button
                    return render_template('success.html',
                                         operation='Compress',
                                         filename=compressed_filename,
                                         pages_count=page_count)

                except ValueError as e:
                    flash(str(e), 'error')
                except Exception as e:
                    app.logger.error(f'Compression error: {e}')
                    record_usage('compress', success=False, error_message=str(e))
                    flash(f'Error compressing PDF: {str(e)}', 'error')
            else:
                flash('Please select a valid PDF file.', 'error')

        return render_template('compress.html', file_uploaded=False)

    @app.route('/rotate', methods=['GET', 'POST'])
    @limiter.limit("30 per hour")
    def rotate_pdf():
        if request.method == 'POST':
            if 'file' in request.files:
                # Check usage limit
                can_proceed, error_msg = check_usage_limit()
                if not can_proceed:
                    flash(error_msg, 'error')
                    return render_template('rotate.html', file_uploaded=False)

                file = request.files['file']
                if file and file.filename and allowed_file(file.filename):
                    try:
                        file_size = validate_file_size(file)

                        safe_filename = sanitize_filename(file.filename)
                        unique_filename = f"{secrets.token_hex(8)}_{safe_filename}"
                        file_path = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)
                        file.save(file_path)

                        try:
                            reader = PdfReader(file_path)
                            page_count = len(reader.pages)
                            if page_count == 0:
                                os.remove(file_path)
                                flash('Invalid PDF: File has no pages.', 'error')
                                return render_template('rotate.html', file_uploaded=False)
                        except Exception as e:
                            if os.path.exists(file_path):
                                os.remove(file_path)
                            app.logger.error(f'PDF validation error: {e}')
                            flash(f'Invalid PDF file: {str(e)}', 'error')
                            return render_template('rotate.html', file_uploaded=False)

                        session['rotate_file'] = unique_filename
                        flash('PDF uploaded successfully!', 'success')
                        return render_template('rotate.html', file_uploaded=True, page_count=page_count, file_name=unique_filename)

                    except ValueError as e:
                        flash(str(e), 'error')
                    except Exception as e:
                        app.logger.error(f'Upload error: {e}')
                        flash(f'Error uploading file: {str(e)}', 'error')
                else:
                    flash('Please select a valid PDF file.', 'error')

            elif 'rotation_angle' in request.form:
                try:
                    file_name = session.get('rotate_file') or request.form.get('file_name')
                    if not file_name:
                        flash('File not found. Please upload again.', 'error')
                        return render_template('rotate.html', file_uploaded=False)

                    file_path = os.path.join(app.config['UPLOAD_FOLDER'], file_name)

                    if not os.path.exists(file_path):
                        flash('File not found. Please upload again.', 'error')
                        session.pop('rotate_file', None)
                        return render_template('rotate.html', file_uploaded=False)

                    rotation_angle = request.form['rotation_angle']
                    apply_to = request.form.get('apply_to', 'all')

                    # Rotate PDF
                    rotated_path = rotate_pdf_pages(file_path, rotation_angle, apply_to)
                    rotated_filename = os.path.basename(rotated_path)

                    # Record usage
                    file_size = os.path.getsize(file_path)
                    reader = PdfReader(file_path)
                    pages_processed = len(reader.pages)
                    record_usage('rotate', file_size=file_size, pages_processed=pages_processed)

                    # Clean up uploaded file
                    try:
                        os.remove(file_path)
                        session.pop('rotate_file', None)
                    except OSError as e:
                        app.logger.error(f'Failed to remove temp file: {e}')

                    # Return success page with download button
                    return render_template('success.html',
                                         operation='Rotate',
                                         filename=rotated_filename,
                                         pages_count=pages_processed)

                except Exception as e:
                    app.logger.error(f'Rotation error: {e}')
                    record_usage('rotate', success=False, error_message=str(e))
                    flash(f'Error rotating PDF: {str(e)}', 'error')
                    session.pop('rotate_file', None)
                    return render_template('rotate.html', file_uploaded=False)

        return render_template('rotate.html', file_uploaded=False)

    @app.route('/merge', methods=['GET', 'POST'])
    @limiter.limit("30 per hour")
    def merge_pdf():
        if request.method == 'POST':
            if 'files' in request.files:
                # Check usage limit
                can_proceed, error_msg = check_usage_limit()
                if not can_proceed:
                    flash(error_msg, 'error')
                    return render_template('merge.html', files_uploaded=False)

                files = request.files.getlist('files')
                max_files = app.config['MAX_MERGE_FILES_PREMIUM'] if current_user.is_premium else app.config['MAX_MERGE_FILES']

                if len(files) < 2:
                    flash('Please select at least 2 PDF files to merge.', 'error')
                    return render_template('merge.html', files_uploaded=False)

                if len(files) > max_files:
                    flash(f'Maximum {max_files} files allowed.', 'error')
                    return render_template('merge.html', files_uploaded=False)

                uploaded_files = []
                try:
                    for file in files:
                        if file and file.filename and allowed_file(file.filename):
                            try:
                                validate_file_size(file)
                                safe_filename = sanitize_filename(file.filename)
                                unique_filename = f"{secrets.token_hex(8)}_{safe_filename}"
                                file_path = os.path.join(app.config['MERGED_FOLDER'], unique_filename)
                                file.save(file_path)
                                uploaded_files.append(unique_filename)
                            except ValueError as e:
                                # Clean up already uploaded files
                                for uploaded_file in uploaded_files:
                                    try:
                                        os.remove(os.path.join(app.config['MERGED_FOLDER'], uploaded_file))
                                    except OSError:
                                        pass
                                flash(str(e), 'error')
                                return render_template('merge.html', files_uploaded=False)

                    if uploaded_files:
                        session['merge_files'] = uploaded_files
                        flash(f'{len(uploaded_files)} files uploaded successfully!', 'success')
                        return render_template('merge.html', files_uploaded=True, file_names=uploaded_files)

                except Exception as e:
                    app.logger.error(f'Merge upload error: {e}')
                    flash(f'Error uploading files: {str(e)}', 'error')

        return render_template('merge.html', files_uploaded=False)

    @app.route('/rearrange', methods=['POST'])
    def rearrange_files():
        """Handle drag-and-drop reordering of files"""
        try:
            file_paths = request.json.get('file_paths', [])
            valid_paths = []
            for path in file_paths:
                if os.path.exists(os.path.join(app.config['MERGED_FOLDER'], path)):
                    valid_paths.append(path)
            return jsonify({'success': True, 'file_paths': valid_paths})
        except Exception as e:
            return jsonify({'success': False, 'error': str(e)}), 400

    @app.route('/merge_files', methods=['POST'])
    def merge_files():
        """Merge multiple PDF files"""
        try:
            file_paths = session.get('merge_files') or request.form.getlist('file_paths[]')

            if not file_paths:
                flash('No files to merge. Please upload files first.', 'error')
                return redirect(url_for('merge_pdf'))

            valid_paths = []
            for file_path in file_paths:
                full_path = os.path.join(app.config['MERGED_FOLDER'], file_path)
                if os.path.exists(full_path):
                    valid_paths.append(file_path)

            if not valid_paths:
                flash('Files not found. Please upload again.', 'error')
                session.pop('merge_files', None)
                return redirect(url_for('merge_pdf'))

            merged_file_path = merge_pdf_files(valid_paths, app.config['MERGED_FOLDER'])
            merged_filename = os.path.basename(merged_file_path)

            # Record usage
            total_size = sum(os.path.getsize(os.path.join(app.config['MERGED_FOLDER'], p)) for p in valid_paths if os.path.exists(os.path.join(app.config['MERGED_FOLDER'], p)))
            record_usage('merge', file_size=total_size, pages_processed=len(valid_paths))

            session.pop('merge_files', None)

            # Return success page with download button
            return render_template('success.html',
                                 operation='Merge',
                                 filename=merged_filename,
                                 files_count=len(valid_paths))

        except Exception as e:
            record_usage('merge', success=False, error_message=str(e))
            flash(f'Error merging PDFs: {str(e)}', 'error')
            session.pop('merge_files', None)
            return redirect(url_for('merge_pdf'))

    @app.route('/pdf2word', methods=['GET', 'POST'])
    @limiter.limit("30 per hour")
    def pdf2word():
        """Convert PDF to Word document"""
        if request.method == 'POST':
            # Check usage limit
            can_proceed, error_msg = check_usage_limit()
            if not can_proceed:
                flash(error_msg, 'error')
                return render_template('pdf2word.html')

            file = request.files.get('file')
            ai_enhanced = request.form.get('ai_enhanced') == 'on'
            
            if file and file.filename and allowed_file(file.filename):
                try:
                    file_size = validate_file_size(file)

                    safe_filename = sanitize_filename(file.filename)
                    unique_filename = f"{secrets.token_hex(8)}_{safe_filename}"
                    file_path = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)
                    file.save(file_path)

                    # Validate PDF
                    try:
                        reader = PdfReader(file_path)
                        page_count = len(reader.pages)
                        if page_count == 0:
                            os.remove(file_path)
                            flash('Invalid PDF: File has no pages.', 'error')
                            return render_template('pdf2word.html')
                    except Exception as e:
                        os.remove(file_path)
                        flash(f'Invalid PDF file: {str(e)}', 'error')
                        return render_template('pdf2word.html')

                    # Check if PDF2Word conversion is available
                    if not PDF2WORD_AVAILABLE:
                        os.remove(file_path)
                        flash('PDF to Word conversion is not available. Please contact support.', 'error')
                        return render_template('pdf2word.html')

                    # AI-Enhanced conversion (Premium)
                    if ai_enhanced:
                        # Check if user has AI credits
                        can_use_ai, ai_error = check_ai_usage_limit()
                        if not can_use_ai:
                            flash(f'AI conversion requires premium access: {ai_error}', 'error')
                            return render_template('pdf2word.html')
                        
                        # Use AI-enhanced conversion
                        word_file_path = convert_pdf_to_word_ai_enhanced(file_path, app.config['SPLIT_FOLDER'])
                        record_usage('pdf2word_ai', file_size=file_size, pages_processed=page_count)
                        
                        # Deduct AI credit
                        if current_user.is_authenticated:
                            current_user.ai_credits = max(0, (current_user.ai_credits or 0) - 1)
                            db.session.commit()
                    else:
                        # Standard conversion
                        word_file_path = convert_pdf_to_word(file_path, app.config['SPLIT_FOLDER'])
                        record_usage('pdf2word', file_size=file_size, pages_processed=page_count)

                    # Clean up uploaded PDF
                    try:
                        os.remove(file_path)
                    except:
                        pass

                    return send_file(word_file_path, as_attachment=True,
                                   download_name=f"converted_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx",
                                   mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document')

                except ValueError as e:
                    flash(str(e), 'error')
                except Exception as e:
                    record_usage('pdf2word', success=False, error_message=str(e))
                    flash(f'Error converting PDF: {str(e)}', 'error')
            else:
                flash('Please select a valid PDF file.', 'error')

        return render_template('pdf2word.html')

    # ========== PDF to JPG ==========

    @app.route('/pdf2jpg', methods=['GET', 'POST'])
    @limiter.limit("30 per hour")
    def pdf2jpg():
        """Convert PDF to JPG images"""
        if request.method == 'POST':
            can_proceed, error_msg = check_usage_limit()
            if not can_proceed:
                flash(error_msg, 'error')
                return render_template('pdf2jpg.html')

            file = request.files.get('file')
            if file and file.filename and allowed_file(file.filename):
                try:
                    file_size = validate_file_size(file)
                    safe_filename = sanitize_filename(file.filename)
                    unique_filename = f"{secrets.token_hex(8)}_{safe_filename}"
                    file_path = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)
                    file.save(file_path)

                    # Convert PDF to images
                    output_path = convert_pdf_to_images(file_path, app.config['SPLIT_FOLDER'])
                    record_usage('pdf2jpg', file_size=file_size)

                    # Clean up
                    try:
                        os.remove(file_path)
                    except:
                        pass

                    return send_file(output_path, as_attachment=True,
                                   download_name=f"images_{datetime.now().strftime('%Y%m%d_%H%M%S')}.zip",
                                   mimetype='application/zip')

                except ValueError as e:
                    flash(str(e), 'error')
                except Exception as e:
                    record_usage('pdf2jpg', success=False, error_message=str(e))
                    flash(f'Error converting PDF: {str(e)}', 'error')
            else:
                flash('Please select a valid PDF file.', 'error')

        return render_template('pdf2jpg.html')

    # ========== PDF to Excel ==========

    @app.route('/pdf2excel', methods=['GET', 'POST'])
    @limiter.limit("30 per hour")
    def pdf2excel():
        """Convert PDF tables to Excel"""
        if request.method == 'POST':
            can_proceed, error_msg = check_usage_limit()
            if not can_proceed:
                flash(error_msg, 'error')
                return render_template('pdf2excel.html')

            file = request.files.get('file')
            if file and file.filename and allowed_file(file.filename):
                try:
                    file_size = validate_file_size(file)
                    safe_filename = sanitize_filename(file.filename)
                    unique_filename = f"{secrets.token_hex(8)}_{safe_filename}"
                    file_path = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)
                    file.save(file_path)

                    # Convert PDF to Excel
                    output_path = convert_pdf_to_excel(file_path, app.config['SPLIT_FOLDER'])
                    record_usage('pdf2excel', file_size=file_size)

                    # Clean up
                    try:
                        os.remove(file_path)
                    except:
                        pass

                    return send_file(output_path, as_attachment=True,
                                   download_name=f"tables_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx",
                                   mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')

                except ValueError as e:
                    flash(str(e), 'error')
                except Exception as e:
                    record_usage('pdf2excel', success=False, error_message=str(e))
                    flash(f'Error converting PDF: {str(e)}', 'error')
            else:
                flash('Please select a valid PDF file.', 'error')

        return render_template('pdf2excel.html')

    # ========== PDF to PowerPoint ==========

    @app.route('/pdf2ppt', methods=['GET', 'POST'])
    @limiter.limit("30 per hour")
    def pdf2ppt():
        """Convert PDF to PowerPoint"""
        if request.method == 'POST':
            can_proceed, error_msg = check_usage_limit()
            if not can_proceed:
                flash(error_msg, 'error')
                return render_template('pdf2ppt.html')

            file = request.files.get('file')
            if file and file.filename and allowed_file(file.filename):
                try:
                    file_size = validate_file_size(file)
                    safe_filename = sanitize_filename(file.filename)
                    unique_filename = f"{secrets.token_hex(8)}_{safe_filename}"
                    file_path = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)
                    file.save(file_path)

                    # Convert PDF to PowerPoint
                    output_path = convert_pdf_to_ppt(file_path, app.config['SPLIT_FOLDER'])
                    record_usage('pdf2ppt', file_size=file_size)

                    # Clean up
                    try:
                        os.remove(file_path)
                    except:
                        pass

                    return send_file(output_path, as_attachment=True,
                                   download_name=f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx",
                                   mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')

                except ValueError as e:
                    flash(str(e), 'error')
                except Exception as e:
                    record_usage('pdf2ppt', success=False, error_message=str(e))
                    flash(f'Error converting PDF: {str(e)}', 'error')
            else:
                flash('Please select a valid PDF file.', 'error')

        return render_template('pdf2ppt.html')

    # ========== OCR (Premium Feature) ==========

    @app.route('/ocr', methods=['GET', 'POST'])
    @limiter.limit("20 per hour")
    def ocr_pdf():
        """OCR - Extract text from scanned PDFs (Premium feature)"""
        # Check if user is logged in
        if not current_user.is_authenticated:
            return render_template('ocr.html', ocr_results=None)
        
        # Check if user is premium
        if not current_user.is_premium:
            return render_template('ocr.html', ocr_results=None)
        
        if request.method == 'POST':
            # Check usage limit
            can_proceed, error_msg = check_usage_limit()
            if not can_proceed:
                flash(error_msg, 'error')
                return render_template('ocr.html', ocr_results=None)
            
            file = request.files.get('file')
            if file and file.filename and allowed_file(file.filename):
                try:
                    file_size = validate_file_size(file)
                    language = request.form.get('language', 'eng')
                    output_format = request.form.get('output_format', 'searchable_pdf')
                    
                    safe_filename = sanitize_filename(file.filename)
                    unique_filename = f"{secrets.token_hex(8)}_{safe_filename}"
                    file_path = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)
                    file.save(file_path)
                    
                    # Validate PDF
                    try:
                        reader = PdfReader(file_path)
                        page_count = len(reader.pages)
                        if page_count == 0:
                            os.remove(file_path)
                            flash('Invalid PDF: File has no pages.', 'error')
                            return render_template('ocr.html', ocr_results=None)
                    except Exception as e:
                        if os.path.exists(file_path):
                            os.remove(file_path)
                        flash(f'Invalid PDF file: {str(e)}', 'error')
                        return render_template('ocr.html', ocr_results=None)
                    
                    # Perform OCR
                    try:
                        from ocr_engine import OCREngine, create_searchable_pdf
                        
                        ocr = OCREngine(dpi=300, language=language)
                        ocr_results = ocr.extract_text_from_pdf(file_path)
                        
                        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                        
                        # Create searchable PDF
                        output_filename = f"ocr_{timestamp}.pdf"
                        output_path = os.path.join(app.config['SPLIT_FOLDER'], output_filename)
                        create_searchable_pdf(file_path, ocr_results, output_path)
                        
                        # Create text file
                        text_filename = f"ocr_text_{timestamp}.txt"
                        text_path = os.path.join(app.config['SPLIT_FOLDER'], text_filename)
                        with open(text_path, 'w', encoding='utf-8') as f:
                            for result in ocr_results:
                                f.write(f"=== Page {result['page']} ===\n")
                                f.write(result['text'])
                                f.write("\n\n")
                        
                        # Record usage
                        record_usage('ocr', file_size=file_size, pages_processed=page_count)
                        
                        # Clean up uploaded file
                        try:
                            os.remove(file_path)
                        except:
                            pass
                        
                        flash('OCR processing completed successfully!', 'success')
                        return render_template('ocr.html', 
                                             ocr_results=ocr_results,
                                             output_filename=output_filename,
                                             text_filename=text_filename)
                        
                    except ImportError as e:
                        app.logger.error(f'OCR import error: {e}')
                        flash('OCR feature is not available. Missing dependencies.', 'error')
                        if os.path.exists(file_path):
                            os.remove(file_path)
                        return render_template('ocr.html', ocr_results=None)
                    except Exception as e:
                        app.logger.error(f'OCR processing error: {e}')
                        record_usage('ocr', success=False, error_message=str(e))
                        flash(f'OCR processing failed: {str(e)}', 'error')
                        if os.path.exists(file_path):
                            os.remove(file_path)
                        return render_template('ocr.html', ocr_results=None)
                    
                except ValueError as e:
                    flash(str(e), 'error')
                except Exception as e:
                    app.logger.error(f'OCR error: {e}')
                    flash(f'Error processing file: {str(e)}', 'error')
            else:
                flash('Please select a valid PDF file.', 'error')
        
        return render_template('ocr.html', ocr_results=None)

    # ========== AI DOCUMENT INTELLIGENCE (Premium Features) ==========

    @app.route('/ai-tools')
    def ai_tools():
        """AI Tools hub page"""
        return render_template('ai-tools.html')

    @app.route('/ai/summarize', methods=['GET', 'POST'])
    @limiter.limit("20 per hour")
    def ai_summarize():
        """AI-powered PDF summarization"""
        if request.method == 'GET':
            return render_template('ai-summarize.html')
        
        # POST - Process the summarization request
        try:
            # Check if AI service is configured
            from ai_services import get_ai_service, summarize_document
            from document_intelligence import process_pdf_for_ai
            
            ai_service = get_ai_service()
            if not ai_service.is_configured():
                return jsonify({
                    'success': False,
                    'error': 'AI service not configured. Please contact administrator.'
                }), 503
            
            # Check AI credits
            can_use_ai, ai_msg = check_ai_usage_limit()
            if not can_use_ai:
                return jsonify({
                    'success': False,
                    'error': ai_msg,
                    'upgrade_required': True
                }), 403
            
            # Get file and options
            file = request.files.get('file')
            length = request.form.get('length', 'medium')
            format_type = request.form.get('format', 'paragraph')
            
            if not file or not file.filename:
                return jsonify({'success': False, 'error': 'No file provided'}), 400
            
            if not allowed_file(file.filename):
                return jsonify({'success': False, 'error': 'Invalid file type. Please upload a PDF.'}), 400
            
            # Save and process file
            try:
                file_size = validate_file_size(file)
            except ValueError as e:
                return jsonify({'success': False, 'error': str(e)}), 400
            
            safe_filename = sanitize_filename(file.filename)
            unique_filename = f"{secrets.token_hex(8)}_{safe_filename}"
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)
            file.save(file_path)
            
            try:
                # Extract text from PDF
                pdf_data = process_pdf_for_ai(file_path)
                
                if not pdf_data['success']:
                    return jsonify({
                        'success': False,
                        'error': pdf_data.get('error', 'Failed to extract text from PDF')
                    }), 400
                
                if not pdf_data['text'] or len(pdf_data['text'].strip()) < 50:
                    return jsonify({
                        'success': False,
                        'error': 'PDF contains too little text to summarize.'
                    }), 400
                
                # Generate summary
                summary_result = summarize_document(
                    pdf_data['text'],
                    length=length,
                    format=format_type
                )
                
                if not summary_result.get('success'):
                    return jsonify({
                        'success': False,
                        'error': summary_result.get('error', 'Summarization failed')
                    }), 500
                
                # Record usage and consume AI credit
                record_usage('ai_summarize', file_size=file_size, pages_processed=pdf_data['page_count'])
                if current_user.is_authenticated:
                    current_user.use_ai_credit()
                
                # Return results
                return jsonify({
                    'success': True,
                    'summary': summary_result.get('summary'),
                    'key_points': summary_result.get('key_points', []),
                    'main_topic': summary_result.get('main_topic'),
                    'document_type': summary_result.get('document_type'),
                    'page_count': pdf_data['page_count'],
                    'original_word_count': pdf_data['word_count'],
                    'summary_word_count': summary_result.get('summary_word_count', 0)
                })
                
            finally:
                # Clean up
                if os.path.exists(file_path):
                    try:
                        os.remove(file_path)
                    except:
                        pass
                        
        except ImportError as e:
            app.logger.error(f'AI module import error: {e}')
            return jsonify({
                'success': False,
                'error': 'AI features are not available. Missing dependencies.'
            }), 503
        except Exception as e:
            app.logger.error(f'AI summarization error: {e}')
            return jsonify({
                'success': False,
                'error': f'An error occurred: {str(e)}'
            }), 500

    @app.route('/ai/extract-tables', methods=['GET', 'POST'])
    @limiter.limit("20 per hour")
    def ai_extract_tables():
        """Extract tables from PDF to CSV/JSON"""
        if request.method == 'GET':
            return render_template('ai-extract-tables.html')
        
        try:
            from document_intelligence import extract_tables_to_format, process_pdf_for_ai
            from ai_services import get_ai_service, detect_tables_in_text
            
            # Check AI credits
            can_use_ai, ai_msg = check_ai_usage_limit()
            if not can_use_ai:
                return jsonify({
                    'success': False,
                    'error': ai_msg,
                    'upgrade_required': True
                }), 403
            
            # Get file and format
            file = request.files.get('file')
            output_format = request.form.get('format', 'csv')
            
            if not file or not file.filename:
                return jsonify({'success': False, 'error': 'No file provided'}), 400
            
            if not allowed_file(file.filename):
                return jsonify({'success': False, 'error': 'Invalid file type. Please upload a PDF.'}), 400
            
            # Save and process file
            try:
                file_size = validate_file_size(file)
            except ValueError as e:
                return jsonify({'success': False, 'error': str(e)}), 400
            
            safe_filename = sanitize_filename(file.filename)
            unique_filename = f"{secrets.token_hex(8)}_{safe_filename}"
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)
            file.save(file_path)
            
            try:
                # Extract tables using pattern matching first
                table_result = extract_tables_to_format(file_path, output_format)
                
                # If no tables found with pattern matching, try AI detection
                if table_result.get('success') and not table_result.get('tables_found'):
                    ai_service = get_ai_service()
                    if ai_service.is_configured():
                        # Get the text first
                        pdf_data = process_pdf_for_ai(file_path)
                        if pdf_data['success'] and pdf_data['text']:
                            ai_tables = detect_tables_in_text(pdf_data['text'])
                            if ai_tables.get('success') and ai_tables.get('tables_found'):
                                table_result = {
                                    'success': True,
                                    'tables_found': True,
                                    'table_count': ai_tables['table_count'],
                                    'tables': ai_tables['tables'],
                                    'formatted_data': json.dumps(ai_tables['tables'], indent=2) if output_format == 'json' else '',
                                    'detection_method': 'ai'
                                }
                
                # Record usage and consume AI credit
                record_usage('ai_extract_tables', file_size=file_size)
                if current_user.is_authenticated:
                    current_user.use_ai_credit()
                
                return jsonify(table_result)
                
            finally:
                # Clean up
                if os.path.exists(file_path):
                    try:
                        os.remove(file_path)
                    except:
                        pass
                        
        except ImportError as e:
            app.logger.error(f'Document intelligence import error: {e}')
            return jsonify({
                'success': False,
                'error': 'Table extraction features are not available.'
            }), 503
        except Exception as e:
            app.logger.error(f'Table extraction error: {e}')
            return jsonify({'success': False, 'error': str(e)}), 500

    @app.route('/ai/process-pdf', methods=['POST'])
    @limiter.limit("30 per hour")
    def ai_process_pdf():
        """Process PDF and extract text for Q&A"""
        try:
            from document_intelligence import process_pdf_for_ai
            
            # Check AI credits
            can_use_ai, ai_msg = check_ai_usage_limit()
            if not can_use_ai:
                return jsonify({
                    'success': False,
                    'error': ai_msg,
                    'upgrade_required': True
                }), 403
            
            file = request.files.get('file')
            if not file or not file.filename:
                return jsonify({'success': False, 'error': 'No file provided'}), 400
            
            if not allowed_file(file.filename):
                return jsonify({'success': False, 'error': 'Invalid file type'}), 400
            
            try:
                file_size = validate_file_size(file)
            except ValueError as e:
                return jsonify({'success': False, 'error': str(e)}), 400
            
            safe_filename = sanitize_filename(file.filename)
            unique_filename = f"{secrets.token_hex(8)}_{safe_filename}"
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)
            file.save(file_path)
            
            try:
                pdf_data = process_pdf_for_ai(file_path)
                
                if not pdf_data['success']:
                    return jsonify({
                        'success': False,
                        'error': pdf_data.get('error', 'Failed to process PDF')
                    }), 400
                
                return jsonify({
                    'success': True,
                    'text': pdf_data['text'],
                    'page_count': pdf_data['page_count'],
                    'word_count': pdf_data['word_count'],
                    'file_name': file.filename
                })
                
            finally:
                if os.path.exists(file_path):
                    try:
                        os.remove(file_path)
                    except:
                        pass
                        
        except ImportError:
            return jsonify({
                'success': False,
                'error': 'Document processing features are not available.'
            }), 503
        except Exception as e:
            app.logger.error(f'PDF processing error: {e}')
            return jsonify({'success': False, 'error': str(e)}), 500

    @app.route('/ai/qa', methods=['GET', 'POST'])
    @limiter.limit("30 per hour")
    def ai_qa():
        """AI-powered PDF Q&A"""
        if request.method == 'GET':
            return render_template('ai-qa.html')
        
        try:
            from ai_services import get_ai_service, answer_question
            
            ai_service = get_ai_service()
            if not ai_service.is_configured():
                return jsonify({
                    'success': False,
                    'error': 'AI service not configured.'
                }), 503
            
            # Check AI credits
            can_use_ai, ai_msg = check_ai_usage_limit()
            if not can_use_ai:
                return jsonify({
                    'success': False,
                    'error': ai_msg,
                    'upgrade_required': True
                }), 403
            
            # Get question and text from request
            data = request.get_json()
            if not data:
                return jsonify({'success': False, 'error': 'No data provided'}), 400
            
            question = data.get('question', '').strip()
            text = data.get('text', '')
            
            if not question:
                return jsonify({'success': False, 'error': 'Please provide a question'}), 400
            
            if not text:
                return jsonify({'success': False, 'error': 'No document text provided'}), 400
            
            # Answer the question
            qa_result = answer_question(text, question)
            
            if not qa_result.get('success'):
                return jsonify({
                    'success': False,
                    'error': qa_result.get('error', 'Failed to answer question')
                }), 500
            
            # Record usage and consume AI credit
            record_usage('ai_qa')
            if current_user.is_authenticated:
                current_user.use_ai_credit()
            
            return jsonify({
                'success': True,
                'answer': qa_result.get('answer'),
                'confidence': qa_result.get('confidence', 'medium'),
                'relevant_quotes': qa_result.get('relevant_quotes', []),
                'question': question
            })
            
        except ImportError:
            return jsonify({
                'success': False,
                'error': 'AI features are not available.'
            }), 503
        except Exception as e:
            app.logger.error(f'AI Q&A error: {e}')
            return jsonify({'success': False, 'error': str(e)}), 500

    # ========== PDF PROCESSING FUNCTIONS ==========

    def split_pdf_pages(pdf_path, pages_to_split, output_dir):
        """Split PDF pages - each page becomes a separate PDF file.
        
        Returns:
            Path to the output file (PDF if single page, ZIP if multiple pages)
        """
        reader = PdfReader(pdf_path)
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        
        # If only one page, return a single PDF directly (no ZIP)
        if len(pages_to_split) == 1:
            page_num = pages_to_split[0]
            if page_num < len(reader.pages):
                writer = PdfWriter()
                writer.add_page(reader.pages[page_num])
                
                output_filename = f"page_{page_num + 1}_{timestamp}.pdf"
                output_path = os.path.join(output_dir, output_filename)
                
                with open(output_path, "wb") as output_pdf:
                    writer.write(output_pdf)
                
                return output_path
        
        # Multiple pages - create a ZIP file
        zip_file_path = os.path.join(output_dir, f"split_pages_{timestamp}.zip")

        with zipfile.ZipFile(zip_file_path, 'w', zipfile.ZIP_DEFLATED) as zip_file:
            for page_num in pages_to_split:
                if page_num >= len(reader.pages):
                    continue

                writer = PdfWriter()
                writer.add_page(reader.pages[page_num])

                split_filename = f"page_{page_num + 1}.pdf"
                split_file_path = os.path.join(output_dir, split_filename)

                with open(split_file_path, "wb") as output_pdf:
                    writer.write(output_pdf)

                zip_file.write(split_file_path, split_filename)

                try:
                    os.remove(split_file_path)
                except:
                    pass

        return zip_file_path

    def split_pdf_by_ranges(pdf_path, ranges_list, output_dir):
        """Split PDF by ranges - each range becomes one PDF document with multiple pages.
        
        Args:
            pdf_path: Path to the source PDF
            ranges_list: List of (start, end) tuples (0-indexed)
            output_dir: Directory to save the output files
            
        Returns:
            Path to the output file (PDF if single range, ZIP if multiple ranges)
        """
        reader = PdfReader(pdf_path)
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        
        # If only one range, return a single PDF directly (no ZIP)
        if len(ranges_list) == 1:
            start, end = ranges_list[0]
            writer = PdfWriter()
            
            # Add all pages in the range
            for page_num in range(start, end + 1):
                if page_num < len(reader.pages):
                    writer.add_page(reader.pages[page_num])
            
            # Create filename based on range
            if start == end:
                output_filename = f"page_{start + 1}_{timestamp}.pdf"
            else:
                output_filename = f"pages_{start + 1}-{end + 1}_{timestamp}.pdf"
            
            output_path = os.path.join(output_dir, output_filename)
            with open(output_path, "wb") as output_pdf:
                writer.write(output_pdf)
            
            return output_path
        
        # Multiple ranges - create a ZIP file
        zip_file_path = os.path.join(output_dir, f"split_ranges_{timestamp}.zip")

        with zipfile.ZipFile(zip_file_path, 'w', zipfile.ZIP_DEFLATED) as zip_file:
            for idx, (start, end) in enumerate(ranges_list):
                writer = PdfWriter()
                
                # Add all pages in the range to this PDF
                for page_num in range(start, end + 1):
                    if page_num < len(reader.pages):
                        writer.add_page(reader.pages[page_num])
                
                # Create filename based on range (convert back to 1-indexed for filename)
                if start == end:
                    split_filename = f"page_{start + 1}.pdf"
                else:
                    split_filename = f"pages_{start + 1}-{end + 1}.pdf"
                
                split_file_path = os.path.join(output_dir, split_filename)
                
                with open(split_file_path, "wb") as output_pdf:
                    writer.write(output_pdf)
                
                zip_file.write(split_file_path, split_filename)
                
                try:
                    os.remove(split_file_path)
                except:
                    pass

        return zip_file_path

    def merge_pdf_files(file_paths, output_dir):
        """Merge multiple PDF files into one"""
        writer = PdfWriter()
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')

        for pdf_file in file_paths:
            pdf_path = os.path.join(output_dir, pdf_file)
            if not os.path.exists(pdf_path):
                continue

            try:
                reader = PdfReader(pdf_path)
                for page in reader.pages:
                    writer.add_page(page)
            except:
                continue

        merged_file_path = os.path.join(output_dir, f"merged_{timestamp}.pdf")
        with open(merged_file_path, "wb") as output_pdf:
            writer.write(output_pdf)

        for pdf_file in file_paths:
            try:
                os.remove(os.path.join(output_dir, pdf_file))
            except:
                pass

        return merged_file_path

    def compress_pdf_file(pdf_path, compression_level='medium'):
        """Compress PDF by removing unnecessary content"""
        reader = PdfReader(pdf_path)
        writer = PdfWriter()
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')

        for page in reader.pages:
            page.compress_content_streams()
            writer.add_page(page)

        compressed_file_path = os.path.join(app.config['SPLIT_FOLDER'], f"compressed_{timestamp}.pdf")
        with open(compressed_file_path, "wb") as output_pdf:
            writer.write(output_pdf)

        return compressed_file_path

    def rotate_pdf_pages(pdf_path, rotation_angle, apply_to='all'):
        """Rotate PDF pages by specified angle"""
        reader = PdfReader(pdf_path)
        writer = PdfWriter()
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')

        angle = int(rotation_angle)

        for idx, page in enumerate(reader.pages):
            should_rotate = apply_to == 'all'
            if apply_to == 'odd' and (idx + 1) % 2 != 0:
                should_rotate = True
            elif apply_to == 'even' and (idx + 1) % 2 == 0:
                should_rotate = True

            if should_rotate:
                page.rotate(angle)

            writer.add_page(page)

        rotated_file_path = os.path.join(app.config['SPLIT_FOLDER'], f"rotated_{timestamp}.pdf")
        with open(rotated_file_path, "wb") as output_pdf:
            writer.write(output_pdf)

        return rotated_file_path

    def convert_pdf_to_word(pdf_path, output_dir):
        """Convert PDF to Word document with tables, images and formatting"""
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        word_file_path = os.path.join(output_dir, f"converted_{timestamp}.docx")

        # Open PDF with PyMuPDF
        pdf_document = fitz.open(pdf_path)
        
        # Create Word document
        doc = Document()
        
        for page_num in range(len(pdf_document)):
            page = pdf_document[page_num]
            
            # Try to extract tables first
            try:
                tables = page.find_tables()
                table_rects = [t.bbox for t in tables] if tables else []
            except:
                tables = None
                table_rects = []
            
            # Extract images
            image_list = page.get_images()
            for img_index, img in enumerate(image_list):
                try:
                    xref = img[0]
                    base_image = pdf_document.extract_image(xref)
                    image_bytes = base_image["image"]
                    
                    # Save image temporarily
                    img_path = os.path.join(output_dir, f"temp_img_{page_num}_{img_index}.png")
                    with open(img_path, "wb") as img_file:
                        img_file.write(image_bytes)
                    
                    # Add to document
                    try:
                        doc.add_picture(img_path, width=Inches(5))
                    except:
                        pass
                    
                    # Clean up temp image
                    try:
                        os.remove(img_path)
                    except:
                        pass
                except:
                    continue
            
            # Add tables
            if tables:
                for table in tables:
                    try:
                        table_data = table.extract()
                        if table_data and len(table_data) > 0:
                            rows = len(table_data)
                            cols = max(len(row) for row in table_data) if table_data else 0
                            if rows > 0 and cols > 0:
                                word_table = doc.add_table(rows=rows, cols=cols)
                                word_table.style = 'Table Grid'
                                for i, row in enumerate(table_data):
                                    for j, cell in enumerate(row):
                                        if j < cols:
                                            word_table.rows[i].cells[j].text = str(cell) if cell else ""
                                doc.add_paragraph("")  # Space after table
                    except:
                        continue
            
            # Extract text with formatting using dict mode
            text_dict = page.get_text("dict")
            
            for block in text_dict.get("blocks", []):
                if block.get("type") == 0:  # Text block
                    # Skip if this block is inside a table
                    block_rect = fitz.Rect(block.get("bbox", (0, 0, 0, 0)))
                    in_table = any(fitz.Rect(tr).intersects(block_rect) for tr in table_rects)
                    if in_table:
                        continue
                    
                    for line in block.get("lines", []):
                        para = doc.add_paragraph()
                        for span in line.get("spans", []):
                            text = span.get("text", "")
                            if text.strip():
                                run = para.add_run(text)
                                
                                # Apply font size
                                font_size = span.get("size", 11)
                                run.font.size = Pt(min(font_size, 72))  # Cap at 72pt
                                
                                # Apply font name
                                font_name = span.get("font", "")
                                if font_name:
                                    # Clean up font name (remove subset prefix like ABCDEF+)
                                    if "+" in font_name:
                                        font_name = font_name.split("+")[-1]
                                    run.font.name = font_name
                                
                                # Apply text color
                                color = span.get("color", 0)
                                if color and color != 0:
                                    # Color is stored as integer, convert to RGB
                                    r = (color >> 16) & 0xFF
                                    g = (color >> 8) & 0xFF
                                    b = color & 0xFF
                                    # Only apply non-black colors
                                    if not (r == 0 and g == 0 and b == 0):
                                        run.font.color.rgb = RGBColor(r, g, b)
                                
                                # Apply bold/italic/etc based on font flags
                                flags = span.get("flags", 0)
                                if flags & 2 ** 0:  # Superscript
                                    run.font.superscript = True
                                if flags & 2 ** 1:  # Italic
                                    run.font.italic = True
                                if flags & 2 ** 2:  # Serifed (skip)
                                    pass
                                if flags & 2 ** 3:  # Monospace
                                    run.font.name = "Courier New"
                                if flags & 2 ** 4:  # Bold
                                    run.font.bold = True
            
            # Add page break between pages (except last)
            if page_num < len(pdf_document) - 1:
                doc.add_page_break()
        
        pdf_document.close()
        doc.save(word_file_path)

        return word_file_path

    def convert_pdf_to_word_ai_enhanced(pdf_path, output_dir):
        """Convert PDF to Word using AI for intelligent document restructuring"""
        from ai_services import enhance_document_for_word
        
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        word_file_path = os.path.join(output_dir, f"ai_converted_{timestamp}.docx")

        # Open PDF with PyMuPDF
        pdf_document = fitz.open(pdf_path)
        
        # Extract all text from PDF
        full_text = ""
        tables_data = []
        
        for page_num in range(len(pdf_document)):
            page = pdf_document[page_num]
            full_text += page.get_text() + "\n\n--- Page Break ---\n\n"
            
            # Extract tables if available
            try:
                tables = page.find_tables()
                if tables:
                    for table in tables:
                        tables_data.append(table.extract())
            except:
                pass
        
        pdf_document.close()
        
        # Send to AI for intelligent restructuring
        ai_result = enhance_document_for_word(full_text, tables_data)
        
        if not ai_result.get('success'):
            # Fallback to standard conversion
            return convert_pdf_to_word(pdf_path, output_dir)
        
        structured_content = ai_result.get('structured_content', {})
        sections = structured_content.get('sections', [])
        
        # Create Word document from AI-structured content
        doc = Document()
        
        # Set document title if available
        title = structured_content.get('title')
        if title:
            heading = doc.add_heading(title, level=0)
            heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Process each section
        for section in sections:
            section_type = section.get('type', 'paragraph')
            
            if section_type == 'heading':
                level = min(section.get('level', 1), 9)
                doc.add_heading(section.get('content', ''), level=level)
                
            elif section_type == 'paragraph':
                content = section.get('content', '')
                para = doc.add_paragraph()
                
                # Parse bold and italic markers
                import re
                parts = re.split(r'(\*\*[^*]+\*\*|\*[^*]+\*)', content)
                for part in parts:
                    if part.startswith('**') and part.endswith('**'):
                        run = para.add_run(part[2:-2])
                        run.bold = True
                    elif part.startswith('*') and part.endswith('*'):
                        run = para.add_run(part[1:-1])
                        run.italic = True
                    else:
                        para.add_run(part)
                        
            elif section_type == 'list':
                style = section.get('style', 'bullet')
                items = section.get('items', [])
                for i, item in enumerate(items):
                    if style == 'numbered':
                        para = doc.add_paragraph(f"{i+1}. {item}")
                    else:
                        para = doc.add_paragraph(item, style='List Bullet')
                        
            elif section_type == 'table':
                headers = section.get('headers', [])
                rows = section.get('rows', [])
                if headers or rows:
                    num_cols = len(headers) if headers else (len(rows[0]) if rows else 0)
                    num_rows = (1 if headers else 0) + len(rows)
                    if num_cols > 0 and num_rows > 0:
                        table = doc.add_table(rows=num_rows, cols=num_cols)
                        table.style = 'Table Grid'
                        
                        row_idx = 0
                        if headers:
                            for j, header in enumerate(headers):
                                cell = table.rows[0].cells[j]
                                cell.text = str(header)
                                # Bold headers
                                for paragraph in cell.paragraphs:
                                    for run in paragraph.runs:
                                        run.bold = True
                            row_idx = 1
                        
                        for row_data in rows:
                            for j, cell_data in enumerate(row_data):
                                if j < num_cols:
                                    table.rows[row_idx].cells[j].text = str(cell_data) if cell_data else ""
                            row_idx += 1
                    doc.add_paragraph("")  # Space after table
                    
            elif section_type == 'quote':
                para = doc.add_paragraph(section.get('content', ''))
                para.style = 'Quote'
                
            elif section_type == 'page_break':
                doc.add_page_break()
        
        doc.save(word_file_path)
        return word_file_path

    def convert_pdf_to_images(pdf_path, output_dir):
        """Convert PDF pages to JPG images and return as ZIP"""
        import zipfile
        from PIL import Image
        
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        zip_path = os.path.join(output_dir, f"images_{timestamp}.zip")
        
        pdf_document = fitz.open(pdf_path)
        
        with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for page_num in range(len(pdf_document)):
                page = pdf_document[page_num]
                # Render at 2x resolution for quality
                mat = fitz.Matrix(2, 2)
                pix = page.get_pixmap(matrix=mat)
                
                img_path = os.path.join(output_dir, f"page_{page_num + 1}.jpg")
                pix.save(img_path)
                
                zipf.write(img_path, f"page_{page_num + 1}.jpg")
                
                # Clean up temp image
                try:
                    os.remove(img_path)
                except:
                    pass
        
        pdf_document.close()
        return zip_path

    def convert_pdf_to_excel(pdf_path, output_dir):
        """Extract tables from PDF and save as Excel"""
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Border, Side
        
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        excel_path = os.path.join(output_dir, f"tables_{timestamp}.xlsx")
        
        pdf_document = fitz.open(pdf_path)
        wb = Workbook()
        ws = wb.active
        ws.title = "Extracted Tables"
        
        # Header style
        header_font = Font(bold=True, color="FFFFFF")
        header_fill = PatternFill(start_color="667EEA", end_color="667EEA", fill_type="solid")
        thin_border = Border(
            left=Side(style='thin'),
            right=Side(style='thin'),
            top=Side(style='thin'),
            bottom=Side(style='thin')
        )
        
        current_row = 1
        table_count = 0
        
        for page_num in range(len(pdf_document)):
            page = pdf_document[page_num]
            
            try:
                tables = page.find_tables()
                if tables:
                    for table in tables:
                        table_data = table.extract()
                        if table_data:
                            table_count += 1
                            
                            # Add table header
                            ws.cell(row=current_row, column=1, value=f"Table {table_count} (Page {page_num + 1})")
                            ws.cell(row=current_row, column=1).font = Font(bold=True, size=12)
                            current_row += 1
                            
                            for i, row in enumerate(table_data):
                                for j, cell in enumerate(row):
                                    cell_obj = ws.cell(row=current_row, column=j + 1, value=str(cell) if cell else "")
                                    cell_obj.border = thin_border
                                    
                                    # Style first row as header
                                    if i == 0:
                                        cell_obj.font = header_font
                                        cell_obj.fill = header_fill
                                current_row += 1
                            
                            current_row += 1  # Space between tables
            except:
                continue
        
        pdf_document.close()
        
        if table_count == 0:
            # If no tables found, extract all text
            ws.cell(row=1, column=1, value="No tables found in PDF. Here is the extracted text:")
            ws.cell(row=1, column=1).font = Font(bold=True)
            
            pdf_document = fitz.open(pdf_path)
            current_row = 3
            for page_num in range(len(pdf_document)):
                page = pdf_document[page_num]
                text = page.get_text()
                for line in text.split('\n'):
                    if line.strip():
                        ws.cell(row=current_row, column=1, value=line.strip())
                        current_row += 1
            pdf_document.close()
        
        wb.save(excel_path)
        return excel_path

    def convert_pdf_to_ppt(pdf_path, output_dir):
        """Convert PDF pages to PowerPoint slides"""
        from pptx import Presentation
        from pptx.util import Inches
        
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        ppt_path = os.path.join(output_dir, f"presentation_{timestamp}.pptx")
        
        pdf_document = fitz.open(pdf_path)
        prs = Presentation()
        
        # Set slide dimensions to match PDF page (approximately)
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        for page_num in range(len(pdf_document)):
            page = pdf_document[page_num]
            
            # Render page as image at high resolution
            mat = fitz.Matrix(2, 2)
            pix = page.get_pixmap(matrix=mat)
            
            img_path = os.path.join(output_dir, f"slide_{page_num + 1}.png")
            pix.save(img_path)
            
            # Add blank slide and insert image
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add image to cover the entire slide
            slide.shapes.add_picture(img_path, Inches(0), Inches(0), 
                                     width=prs.slide_width, height=prs.slide_height)
            
            # Clean up temp image
            try:
                os.remove(img_path)
            except:
                pass
        
        pdf_document.close()
        prs.save(ppt_path)
        return ppt_path

    return app


# ========== APPLICATION INSTANCE ==========

app = create_app()

if __name__ == "__main__":
    port = int(os.environ.get('PORT', 5000))
    host = os.environ.get('HOST', '0.0.0.0')
    debug = os.environ.get('FLASK_DEBUG', 'False').lower() == 'true'

    with app.app_context():
        db.create_all()  # Create tables if they don't exist

    app.run(host=host, port=port, debug=debug)
